import streamlit as st
import pandas as pd
import os
import json
import re
import nltk
from PyPDF2 import PdfReader
from docx import Document
import pptx
import openpyxl
from nltk.corpus import stopwords
from langchain_openai import ChatOpenAI
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import joblib
import numpy as np
from sentence_transformers import SentenceTransformer, util

nltk.download("stopwords")

############################################
# 1. EXTRACTION TEXTE POUR TOUS FORMATS
############################################

def extraire_texte_fichier(uploaded_file):
    extension = uploaded_file.name.lower().split(".")[-1]

    if extension == "pdf":
        reader = PdfReader(uploaded_file)
        texte = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                texte += t + "\n"
        return texte

    elif extension == "docx":
        doc = Document(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs])

    elif extension == "pptx":
        prs = pptx.Presentation(uploaded_file)
        texte = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texte += shape.text + "\n"
        return texte

    elif extension in ["xlsx", "xls"]:
        wb = openpyxl.load_workbook(uploaded_file, data_only=True)
        texte = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                ligne = " ".join([str(cell) for cell in row if cell is not None])
                texte += ligne + "\n"
        return texte

    else:
        raise ValueError("Format non supporté : " + extension)


############################################
# 2. ANALYSE PAR LLM (TOUS FORMATS)
############################################

def analyser_document_avec_llm(uploaded_file, llm):

    texte = extraire_texte_fichier(uploaded_file)

    def decouper_texte(texte, taille=2500):
        return [texte[i:i+taille] for i in range(0, len(texte), taille)]

    chunks = decouper_texte(texte)
    analyses = []

    for chunk in chunks:
        prompt = f"""
Tu es un expert en analyse de documents cinématographiques.

Retourne STRICTEMENT un JSON valide, sans texte avant ou après.
PAS de ```.

Champs :
- titre : string
- mots_cles : liste de strings
- description : string
- categorie : "science-fiction" | "cinéma" | "policier"

Document :
\"\"\"
{chunk}
\"\"\"
"""
        response = llm.invoke(prompt)
        raw = response.content

        match = re.search(r"\{.*?\}", raw, flags=re.DOTALL)
        if match:
            try:
                analyses.append(json.loads(match.group(0)))
            except:
                pass

    if not analyses:
        raise ValueError("Aucun JSON valide trouvé dans le document")

    fusion = {
        "titre": analyses[0]["titre"],
        "mots_cles": list({mot for a in analyses for mot in a["mots_cles"]}),
        "description": " ".join(a["description"] for a in analyses),
        "categorie": analyses[0]["categorie"]
    }

    return fusion


############################################
# 3. ANALYSE COMPLÈTE (LLM + ML)
############################################

def analyser_document_complet(uploaded_file, llm, clf, vectorizer):

    analyse = analyser_document_avec_llm(uploaded_file, llm)

    titre = analyse["titre"]
    description = analyse["description"]
    mots_cles = ", ".join(analyse["mots_cles"])
    categorie_llm = analyse["categorie"]

    embedding = vectorizer.transform([description])
    probas = clf.predict_proba(embedding)[0]

    categorie_ml = clf.classes_[probas.argmax()]
    confiance_ml = probas.max()

    if confiance_ml < 0.7:
        categorie_finale = categorie_llm
        source = "LLM"
    else:
        categorie_finale = categorie_ml
        source = "ML"

    return {
        "titre": titre,
        "description": description,
        "mots_cles": mots_cles,
        "categorie_llm": categorie_llm,
        "categorie_ml": categorie_ml,
        "confiance_ml": float(confiance_ml),
        "categorie_finale": categorie_finale,
        "source_decision": source
    }


############################################
# 4. AJOUT AU DATAFRAME
############################################

def ajouter_document_au_dataframe(uploaded_file, llm, dataframe_path="documents_analyzed_B20.csv"):

    if os.path.exists(dataframe_path):
        df = pd.read_csv(dataframe_path)
    else:
        df = pd.DataFrame(columns=[
            "file_path", "titre", "mots_cles", "description",
            "categorie_llm", "categorie_ml", "confiance_ml",
            "categorie_finale", "source_decision"
        ])

    file_name = uploaded_file.name

    if file_name in df["file_path"].values:
        st.warning("⚠️ Ce fichier a déjà été analysé")
        return df

    analyse = analyser_document_complet(uploaded_file, llm, clf, vectorizer)

    nouvelle_ligne = {
        "file_path": file_name,
        "titre": analyse["titre"],
        "mots_cles": analyse["mots_cles"],
        "description": analyse["description"],
        "categorie_llm": analyse["categorie_llm"],
        "categorie_ml": analyse["categorie_ml"],
        "confiance_ml": analyse["confiance_ml"],
        "categorie_finale": analyse["categorie_finale"],
        "source_decision": analyse["source_decision"]
    }

    df = pd.concat([df, pd.DataFrame([nouvelle_ligne])], ignore_index=True)
    df.to_csv(dataframe_path, index=False)

    return df


############################################
# 5. RECHERCHE SÉMANTIQUE
############################################

def recherche_semantique(df, requete, top_k=2):

    corpus = (df["titre"].fillna("") + " " +
              df["mots_cles"].fillna("") + " " +
              df["description"].fillna("")).tolist()

    vectorizer = TfidfVectorizer(stop_words="english")
    tfidf_matrix = vectorizer.fit_transform(corpus)

    query_vec = vectorizer.transform([requete])
    scores = cosine_similarity(query_vec, tfidf_matrix).flatten()

    top_indices = scores.argsort()[::-1][:top_k]
    resultats = df.iloc[top_indices].copy()
    resultats["score"] = scores[top_indices]

    return resultats


############################################
# 6. FILTRES
############################################

MOTS_CLES_FIXES = [
    "robotique",
    "IA",
    "automatisation",
    "technologie",
    "innovation",
    "futur",
    "espace",
    "astronautes",
    "trou de ver",
    "exploration",
    "galaxie",
    "mythologie",
    "légendes",
    "symbolisme",
    "esthétique",
    "art",
    "peinture",
    "sculpture",
    "industrie",
    "histoire",
    "psychologie",
    "enquête",
    "scène de crime",
    "manipulation",
    "tension",
    "mystère",
    "récit",
    "philosophie",
    "éthique",
    "pouvoir",
    "identité",
    "énergie",
    "Force",
    "Jedi",
    "Sith",
    "économie",
    "politique",
    "société",
    "culture",
    "créatures",
    "mythes",
    "univers",
    "cosmos",
    "temps",
    "espace-temps",
    "réalité",
    "rêves",
    "mémoire",
    "perception"
]

CATEGORIES_FIXES = [
    "science-fiction",
    "cinéma",
    "policier"
]

def rechercher_par_categorie(df, categorie):
    return df[df["categorie_finale"] == categorie]

def filtrer_par_mot_cle(df, mot_cle):
    return df[df["mots_cles"].str.contains(mot_cle, case=False, na=False)]

def filtrer_par_mots_cles(df, mots):
    pattern = "|".join(mots)

    return df[
        df["titre"].str.contains(pattern, case=False, na=False)
        | df["mots_cles"].str.contains(pattern, case=False, na=False)
        | df["description"].str.contains(pattern, case=False, na=False)
        | df["categorie_finale"].str.contains(pattern, case=False, na=False)
    ]

def filtrage_combine(df, categorie=None, tags=None):
    result = df.copy()

    # 1) Filtre catégorie si fournie
    if categorie:
        result = result[result["categorie_finale"].str.lower() == categorie.lower()]

    # Si aucune catégorie trouvée → on retourne directement (normal)
    if result.empty:
        return result

    # 2) Filtre mots-clés si fournis
    if tags:
        final_result = filtrer_par_mots_cles(result, tags)

        # Si aucun résultat → on garde seulement le filtre catégorie
        if final_result.empty:
            return result
        else:
            return final_result

    # Si pas de tags → juste le filtre catégorie
    return result

import pandas as pd
import re

def recherche_lexicale(df, requete):
    """
    Recherche lexicale dans toutes les colonnes textuelles du DataFrame.
    
    df : DataFrame Pandas
    requete : string (phrase libre)
    """

    # Si requête vide → renvoyer tout
    if not requete or not isinstance(requete, str):
        return df

    # 1) Tokenisation simple
    tokens = re.findall(r"\w+", requete.lower())

    # 2) Stopwords basiques
    stopwords = {"le","la","les","un","une","des","de","du","dans","pour","avec","et","ou","à","au"}
    tokens = [t for t in tokens if t not in stopwords]

    if not tokens:
        return df

    # 3) Regex OR
    pattern = "|".join(map(re.escape, tokens))

    # 4) Sélectionner uniquement les colonnes textuelles
    colonnes_textuelles = df.select_dtypes(include=["object", "string"]).columns

    # 5) Appliquer la recherche sur toutes les colonnes textuelles
    masque = df[colonnes_textuelles].apply(
        lambda row: row.astype(str).str.contains(pattern, case=False, regex=True).any(),
        axis=1
    )

    return df[masque]


############################################
# 7. INTERFACE STREAMLIT
############################################

st.title("🔎 Analyse de documents (PDF, Word, Excel, PowerPoint)")

OPENAI_BASE_URL="http://localhost:11434/v1"
OPENAI_API_KEY="ollama"
GENERATOR_MODEL_NAME = "llama3:latest"

llm = ChatOpenAI(
    model=GENERATOR_MODEL_NAME,
    temperature=0,
    max_tokens=None,
    base_url=OPENAI_BASE_URL,
    api_key=OPENAI_API_KEY,
)

clf = joblib.load("modele_ml_lr.joblib")
vectorizer = joblib.load("vectorizer2.joblib")

# Chargement de la base existante
CSV_PATH = "documents_analyzed_B20.csv"

if not os.path.exists(CSV_PATH):
    st.error(f"❌ La base de données '{CSV_PATH}' est introuvable.")
    st.stop()

df = pd.read_csv(CSV_PATH)

types_acceptes = ["pdf", "docx", "xlsx", "xls", "pptx"]

menu = st.sidebar.radio(
    "Choix de l'action",
    ["Analyser un document", "Ajouter au DataFrame", "Rechercher"]
)

############################################
# 8. ACTIONS STREAMLIT
############################################

if menu == "Analyser un document":
    st.header("Analyse du document par LLM + ML")
    fichier = st.file_uploader("Sélection du fichier", type=types_acceptes)

    if fichier:
        resultat = analyser_document_complet(fichier, llm, clf, vectorizer)
        st.json(resultat)

elif menu == "Ajouter au DataFrame":
    st.header("Ajouter un document au DataFrame")
    fichier = st.file_uploader("Sélection du fichier", type=types_acceptes)

    if fichier:
        df = ajouter_document_au_dataframe(fichier, llm)
        st.success("Document analysé et ajouté au DataFrame")
        st.dataframe(df)

elif menu == "Rechercher":
    st.header("Recherche")

    if os.path.exists(CSV_PATH):
        df = pd.read_csv(CSV_PATH)
        st.dataframe(df)

        choix = st.radio("Type de recherche", ["Recherche combinée", "Recherche lexicale", "Recherche sémantique"])

        if choix == "Recherche sémantique":
            st.header("Recherche sémantique")
            requete = st.text_input("Recherche sémantique :")
            top_k = st.slider("Nombre de résultats", min_value=1, max_value=10, value=3)
            if requete:
                st.dataframe(recherche_semantique(df, requete, top_k=top_k))
        
        elif choix == "Recherche combinée":
            st.header("Recherche combinée")

            # Catégorie
            categorie = st.selectbox(
                "Catégorie :",
                [""] + CATEGORIES_FIXES
            )
            categorie = categorie if categorie != "" else None

            # Tags
            tags = st.multiselect(
                "Mots-clés:",
                options=MOTS_CLES_FIXES
            )

            if st.button("Lancer la recherche"):
                resultats = filtrage_combine(df, categorie, tags)
                st.dataframe(resultats)

        elif choix == "Recherche lexicale":
            st.header("Recherche lexicale")
            requete = st.text_input("Recherche lexicale :")
            if requete :
                st.dataframe(recherche_lexicale(df, requete))
    else:
        st.warning("Base de données vide : ajoute d'abord des documents.")





